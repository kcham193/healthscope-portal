"""
Build HealthScope_Technical_Deck.pptx — technical briefing deck for data
engineers / architects.

Narrative arc (per directive):
    Introduction  →  Component-by-component (what + how)  →  Finalize

Reuses the design tokens and helper patterns from
build_healthscope_presentation.py in hf-data-portal (navy + teal + amber
on white; typography-led; thin dividers) but goes deeper on architecture,
method trade-offs, and real numbers.

Deps:
    pip install python-pptx

Usage (from the deploy clone root):
    python scripts/presentation/build_tech_slides.py

Output:
    HealthScope_Technical_Deck.pptx  (in the current working directory)

Slides:

    Introduction
     1. Title
     2. What is HealthScope
     3. The gap MoH registries leave
     4. Overview — one foundation, four components, one assistant

    Data foundation
     5. Foundation — what
     6. Foundation — how

    Component 1 — Registry Explorer
     7. Registry Explorer — what
     8. Registry Explorer — how

    Component 2 — Population and Access
     9. Population and Access — what
    10. Population and Access — how (methodology)
    11. Population and Access — how (Tanzania R pipeline)

    Component 3 — Service Intelligence
    12. Service Intelligence — what
    13. Service Intelligence — how

    Component 4 — Find Facility
    14. Find Facility — what
    15. Find Facility — how

    Component 5 — Mapy chatbot
    16. Mapy — what
    17. Mapy — how

    Finalize
    18. Where we are and where we're heading
    19. Team + sources + links
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ─── design tokens ─────────────────────────────────────────────────────────
NAVY       = RGBColor(0x1A, 0x3C, 0x5E)
TEAL       = RGBColor(0x3E, 0x7C, 0xB1)
INK        = RGBColor(0x1F, 0x2A, 0x37)
GREY       = RGBColor(0x6B, 0x72, 0x80)
LIGHT_GREY = RGBColor(0xCB, 0xD0, 0xD6)
PALE       = RGBColor(0xF3, 0xF5, 0xF8)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
AMBER      = RGBColor(0xC9, 0xA2, 0x27)  # HealthScope brand gold
GREEN_1    = RGBColor(0x1A, 0x98, 0x50)
GREEN_2    = RGBColor(0xA6, 0xD9, 0x6A)
ORANGE     = RGBColor(0xFD, 0xAE, 0x61)
RED        = RGBColor(0xD7, 0x30, 0x27)
PURPLE     = RGBColor(0x54, 0x27, 0x8F)

TITLE_FONT = "Calibri"
BODY_FONT  = "Calibri"
CODE_FONT  = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ─── presentation setup ────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK_LAYOUT = prs.slide_layouts[6]


def new_slide():
    s = prs.slides.add_slide(BLANK_LAYOUT)
    force_white_background(s)
    return s


def force_white_background(slide):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE


# ─── low-level helpers ─────────────────────────────────────────────────────
def add_text(slide, x, y, w, h, text, *, font=BODY_FONT, size=14, bold=False,
             italic=False, color=INK, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.25):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=15, color=INK,
                line_spacing=1.35, bullet_color=AMBER):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        if i > 0:
            p.space_before = Pt(4)
        r_dot = p.add_run()
        r_dot.text = "• "
        r_dot.font.name = BODY_FONT
        r_dot.font.size = Pt(size)
        r_dot.font.bold = True
        r_dot.font.color.rgb = bullet_color

        r_body = p.add_run()
        r_body.text = item
        r_body.font.name = BODY_FONT
        r_body.font.size = Pt(size)
        r_body.font.color.rgb = color
    return tb


def add_line(slide, x1, y1, x2, y2, *, color=NAVY, weight=1.0):
    ln = slide.shapes.add_connector(1, x1, y1, x2, y2)
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln


def add_rect(slide, x, y, w, h, *, fill=None, line_color=LIGHT_GREY,
             line_weight=0.75, no_line=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if no_line:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_weight)
    return shp


def add_code_block(slide, x, y, w, h, code, *, size=11):
    add_rect(slide, x, y, w, h, fill=PALE, no_line=True)
    tb = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.12),
                                  w - Inches(0.3), h - Inches(0.24))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, ln in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.15
        r = p.add_run()
        r.text = ln if ln else " "
        r.font.name = CODE_FONT
        r.font.size = Pt(size)
        r.font.color.rgb = INK
    return tb


def add_header(slide, title, *, eyebrow=None, subtitle=None):
    if eyebrow:
        add_text(slide, Inches(0.7), Inches(0.55), Inches(11), Inches(0.3),
                 eyebrow.upper(), size=10, bold=True, color=TEAL,
                 line_spacing=1.0)
    add_text(slide, Inches(0.7), Inches(0.85), Inches(12), Inches(0.95),
             title, font=TITLE_FONT, size=30, bold=True, color=NAVY,
             line_spacing=1.05)
    if subtitle:
        add_text(slide, Inches(0.7), Inches(1.62), Inches(12), Inches(0.55),
                 subtitle, size=14, color=GREY, italic=True, line_spacing=1.25)
        rule_y = Inches(2.1)
    else:
        rule_y = Inches(1.72)
    add_line(slide, Inches(0.7), rule_y, Inches(1.9), rule_y,
             color=NAVY, weight=1.5)


def add_footer(slide, page_number, total):
    add_line(slide, Inches(0.7), Inches(7.15), Inches(12.6), Inches(7.15),
             color=LIGHT_GREY, weight=0.5)
    add_text(slide, Inches(0.7), Inches(7.22), Inches(9), Inches(0.25),
             "HealthScope  ·  Technical build overview  ·  Malaria Atlas Project — Dar es Salaam Node",
             size=9, color=GREY)
    add_text(slide, Inches(11.5), Inches(7.22), Inches(1.1), Inches(0.25),
             f"{page_number:02d} / {total:02d}",
             size=9, color=GREY, align=PP_ALIGN.RIGHT)


def add_section_divider_pill(slide, x, y, label, *, fill=AMBER):
    """Small horizontal pill used as a section marker under the header."""
    add_rect(slide, x, y, Inches(1.6), Inches(0.3), fill=fill, no_line=True)
    add_text(slide, x, y + Inches(0.02), Inches(1.6), Inches(0.28),
             label.upper(), size=9, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, line_spacing=1.0)


def add_what_how_badge(slide, x, y, kind):
    """kind = 'what' or 'how' — coloured tag."""
    color = TEAL if kind == "what" else AMBER
    label = kind.upper()
    add_rect(slide, x, y, Inches(0.9), Inches(0.28), fill=color, no_line=True)
    add_text(slide, x, y + Inches(0.02), Inches(0.9), Inches(0.26),
             label, size=9, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, line_spacing=1.0)


# ═══════════════════════════════════════════════════════════════════════════
# INTRODUCTION
# ═══════════════════════════════════════════════════════════════════════════

def slide_title(total):
    s = new_slide()
    add_rect(s, Emu(0), Emu(0), Inches(4.5), SLIDE_H, fill=NAVY, no_line=True)
    add_rect(s, Inches(0.5), Inches(2.7), Inches(0.08), Inches(0.9),
             fill=AMBER, no_line=True)
    add_text(s, Inches(0.75), Inches(2.6), Inches(3.5), Inches(0.5),
             "HEALTHSCOPE", font=TITLE_FONT, size=13, bold=True, color=AMBER,
             line_spacing=1.0)
    add_text(s, Inches(0.75), Inches(3.05), Inches(3.5), Inches(1.5),
             "Health Facility\nIntelligence Portal", font=TITLE_FONT,
             size=32, bold=True, color=WHITE, line_spacing=1.05)
    add_text(s, Inches(0.75), Inches(5.15), Inches(3.5), Inches(1.0),
             "Malaria Atlas Project\nDar es Salaam Node",
             size=11, color=RGBColor(0xC5, 0xD1, 0xDE), line_spacing=1.4)

    add_text(s, Inches(5.3), Inches(2.5), Inches(7.5), Inches(0.4),
             "TECHNICAL BUILD OVERVIEW", size=11, bold=True, color=TEAL,
             line_spacing=1.0)
    add_text(s, Inches(5.3), Inches(2.9), Inches(7.5), Inches(2.2),
             "How we built HealthScope:\nthe foundation, each component,\nand how they fit together.",
             font=TITLE_FONT, size=28, bold=True, color=NAVY,
             line_spacing=1.15)
    add_line(s, Inches(5.3), Inches(5.15), Inches(6.5), Inches(5.15),
             color=NAVY, weight=1.5)
    add_text(s, Inches(5.3), Inches(5.25), Inches(7.5), Inches(1.2),
             "From MoH facility registries to a public, reproducible, "
             "open-source portal.",
             size=15, color=GREY, italic=True, line_spacing=1.45)


def slide_what_is_healthscope(total, page):
    s = new_slide()
    add_header(s, "An open portal for health-facility access analytics",
               eyebrow="What is HealthScope",
               subtitle="One-line brief: MoH facility data + population + reachability, in a single public website")

    add_bullets(s, Inches(0.7), Inches(2.55), Inches(7.0), Inches(3.7), [
        "Open-access web portal — free, no login, hosted on GitHub Pages.",
        "Direct-from-MoH facility data across 8 malaria-endemic SSA countries, "
        "standardised to a shared schema.",
        "Answers not just 'where facilities exist', but 'who they serve' and "
        "'whether people can realistically reach them'.",
        "Built by the Malaria Atlas Project, Dar es Salaam Node — data engineers, "
        "GIS specialists, and epidemiologists.",
        "Fully reproducible: every table, chart, and map comes from a versioned "
        "ETL pipeline you can rerun from source.",
    ], size=14)

    # Right column: headline stats
    x = Inches(8.3)
    add_text(s, x, Inches(2.55), Inches(4.3), Inches(0.4),
             "AT A GLANCE", size=11, bold=True, color=TEAL, line_spacing=1.0)
    stats = [
        ("8",         "Countries live"),
        ("~131,700",  "Facilities"),
        ("100 m",     "Analysis resolution"),
        ("0",         "Backend servers"),
    ]
    row_y = Inches(2.95)
    for i, (v, l) in enumerate(stats):
        y = row_y + Inches(0.85) * i
        add_rect(s, x, y, Inches(4.3), Inches(0.72), fill=PALE, no_line=True)
        add_text(s, x + Inches(0.25), y + Inches(0.1), Inches(2.2),
                 Inches(0.55), v, font=TITLE_FONT, size=22, bold=True,
                 color=AMBER, line_spacing=1.0)
        add_text(s, x + Inches(2.4), y + Inches(0.18), Inches(1.8),
                 Inches(0.5), l, size=11, color=INK, line_spacing=1.2)
    add_footer(s, page, total)


def slide_the_gap(total, page):
    s = new_slide()
    add_header(s, "MoH registries answer 'where', not 'who reaches it'",
               eyebrow="The gap",
               subtitle="What we can and can't do with a bare MoH facility list")

    add_bullets(s, Inches(0.7), Inches(2.55), Inches(6.5), Inches(3.5), [
        "MoH registries publish facility name + coordinates + basic metadata.",
        "They rarely publish population served, road access, or per-service availability.",
        "Static PDF maps or spreadsheets don't inform access planning.",
        "'Nearest facility' is a straight-line concept in most tools — even when "
        "the real question is 'can I get there in an hour?'",
    ])

    add_text(s, Inches(7.7), Inches(2.55), Inches(5.0), Inches(0.3),
             "WHAT RAW MoH DATA LOOKS LIKE",
             size=11, bold=True, color=TEAL, line_spacing=1.0)
    add_rect(s, Inches(7.7), Inches(2.9), Inches(5.0), Inches(3.5),
             fill=PALE, no_line=True)
    add_text(s, Inches(7.85), Inches(3.0), Inches(4.7), Inches(3.3),
             "facility_code   100675-8\n"
             "facility_name   CANOSSA DISPENSARY\n"
             "admin1          Arusha\n"
             "admin2          Arusha\n"
             "facility_type   Dispensary\n"
             "ownership       Private (FBO)\n"
             "latitude        -3.45802\n"
             "longitude       36.71262\n"
             "\n"
             "# What's missing:\n"
             "#   catchment_population       ?\n"
             "#   pop_within_5km             ?\n"
             "#   pop_within_60min_travel    ?\n"
             "#   is_reachable_by_road       ?",
             font=CODE_FONT, size=10, color=INK, line_spacing=1.35)
    add_footer(s, page, total)


def slide_overview(total, page):
    s = new_slide()
    add_header(s, "One foundation. Four components. One assistant.",
               eyebrow="What we built",
               subtitle="Everything sits on the standardised registry lake at the bottom")

    # Foundation strip
    fy = Inches(2.55)
    add_rect(s, Inches(0.7), fy, Inches(11.9), Inches(0.9), fill=NAVY,
             no_line=True)
    add_text(s, Inches(0.9), fy + Inches(0.15), Inches(11.5), Inches(0.35),
             "FOUNDATION", size=10, bold=True, color=AMBER, line_spacing=1.0)
    add_text(s, Inches(0.9), fy + Inches(0.42), Inches(11.5), Inches(0.5),
             "Data foundation — MoH-sourced facility CSVs, one canonical schema, "
             "8 countries, ~131,700 rows.",
             size=14, color=WHITE, line_spacing=1.2)

    # 4 component boxes
    box_y = Inches(3.75)
    box_w = Inches(2.9)
    box_h = Inches(2.35)
    gap = Inches(0.15)
    total_w = 4 * box_w + 3 * gap
    start_x = (SLIDE_W - total_w) / 2
    components = [
        ("1", "Registry\nExplorer",
         "Coverage map,\ninteractive filters,\nCSV export.",
         AMBER),
        ("2", "Population\n& Access",
         "Voronoi catchments\n+ travel-time bands\n+ district rankings.",
         GREEN_1),
        ("3", "Service\nIntelligence",
         "7-domain service\ncoverage across\nthe registry.",
         ORANGE),
        ("4", "Find Facility",
         "Public nearest-\nfacility search with\nreal road distance.",
         TEAL),
    ]
    for i, (num, title, body, color) in enumerate(components):
        x = start_x + i * (box_w + gap)
        add_rect(s, x, box_y, box_w, box_h, fill=WHITE,
                 line_color=LIGHT_GREY, line_weight=0.75)
        add_rect(s, x, box_y, box_w, Inches(0.07), fill=color, no_line=True)
        add_text(s, x + Inches(0.2), box_y + Inches(0.2), Inches(0.5),
                 Inches(0.4), num, font=TITLE_FONT, size=22, bold=True,
                 color=color, line_spacing=1.0)
        add_text(s, x + Inches(0.2), box_y + Inches(0.75), box_w - Inches(0.4),
                 Inches(0.7), title, font=TITLE_FONT, size=13, bold=True,
                 color=NAVY, line_spacing=1.15)
        add_text(s, x + Inches(0.2), box_y + Inches(1.45), box_w - Inches(0.4),
                 Inches(0.9), body, size=10, color=INK, line_spacing=1.35)

    # Assistant strip
    ay = Inches(6.3)
    add_rect(s, Inches(0.7), ay, Inches(11.9), Inches(0.65), fill=PALE,
             no_line=True)
    add_text(s, Inches(0.9), ay + Inches(0.1), Inches(11.5), Inches(0.35),
             "ASSISTANT — MAPY", size=10, bold=True, color=AMBER,
             line_spacing=1.0)
    add_text(s, Inches(0.9), ay + Inches(0.35), Inches(11.5), Inches(0.35),
             "A scoped chatbot on every page. Helps visitors find pages and "
             "explains MAP context. Refuses out-of-scope questions.",
             size=11, color=INK, line_spacing=1.2)
    add_footer(s, page, total)


# ═══════════════════════════════════════════════════════════════════════════
# DATA FOUNDATION
# ═══════════════════════════════════════════════════════════════════════════

def slide_foundation_what(total, page):
    s = new_slide()
    add_what_how_badge(s, Inches(0.7), Inches(0.4), "what")
    add_header(s, "Direct-from-MoH facility data, unified across 8 countries",
               eyebrow="Data foundation — what",
               subtitle="One canonical schema; every downstream analysis reads from these files")

    add_bullets(s, Inches(0.7), Inches(2.55), Inches(7.5), Inches(3.7), [
        "Primary sources are MoH facility registries — no second-hand aggregators.",
        "Every country is harmonised to one 30-column CSV per country.",
        "Coordinate-validated (against country bounding box) and stable-UID'd.",
        "Portal pages consume the CSVs directly via Papa Parse in the browser — "
        "no database, no API layer.",
    ], size=14)

    # Right column: per-country facility counts + sources
    x = Inches(8.6)
    add_text(s, x, Inches(2.55), Inches(4.2), Inches(0.4),
             "COUNTRIES × SOURCES",
             size=11, bold=True, color=TEAL, line_spacing=1.0)
    rows = [
        ("Botswana", "1,076",  "MFL"),
        ("Ethiopia", "40,035", "MFR v2"),
        ("Kenya",    "17,353", "KMHFR"),
        ("Malawi",   "1,929",  "MHFR"),
        ("Nigeria",  "51,023", "Nigeria HFR"),
        ("Tanzania", "13,075", "HFR Portal"),
        ("Uganda",   "8,512",  "NHFR"),
        ("Zambia",   "3,731",  "MFL"),
    ]
    row_y = Inches(2.95)
    row_h = Inches(0.36)
    for i, (name, count, src) in enumerate(rows):
        y = row_y + row_h * i
        if i % 2 == 0:
            add_rect(s, x, y, Inches(4.2), row_h, fill=PALE, no_line=True)
        add_text(s, x + Inches(0.1), y + Inches(0.05), Inches(1.3),
                 row_h, name, size=11, color=INK, line_spacing=1.0)
        add_text(s, x + Inches(1.5), y + Inches(0.05), Inches(1.1),
                 row_h, count, size=11, bold=True, color=AMBER,
                 align=PP_ALIGN.RIGHT, line_spacing=1.0)
        add_text(s, x + Inches(2.7), y + Inches(0.05), Inches(1.4),
                 row_h, src, size=10, color=GREY, italic=True,
                 line_spacing=1.0)
    add_footer(s, page, total)


def slide_foundation_how(total, page):
    s = new_slide()
    add_what_how_badge(s, Inches(0.7), Inches(0.4), "how")
    add_header(s, "ETL pipeline v2 — extract, transform, load",
               eyebrow="Data foundation — how",
               subtitle="Idempotent, versionable, per-country. Reruns take minutes.")

    # 3-step flow
    step_y = Inches(2.55)
    step_h = Inches(1.55)
    step_w = Inches(3.7)
    gap = Inches(0.35)
    total_w = 3 * step_w + 2 * gap
    start_x = (SLIDE_W - total_w) / 2
    steps = [
        ("EXTRACT",
         "Fetch the latest per-country registry\n"
         "(CSV export, DHIS2 API, MFR portal).",
         AMBER),
        ("TRANSFORM",
         "Standardise columns, validate coords\n"
         "against bbox, mint stable UIDs, map\n"
         "raw services onto the 7-domain hierarchy.",
         TEAL),
        ("LOAD",
         "Publish versioned per-country CSVs\n"
         "under etl/data/processed/. Portal reads\n"
         "them at build time via Quarto resources.",
         NAVY),
    ]
    for i, (name, body, color) in enumerate(steps):
        x = start_x + i * (step_w + gap)
        add_rect(s, x, step_y, step_w, step_h, fill=WHITE, line_color=color,
                 line_weight=1.25)
        add_text(s, x + Inches(0.2), step_y + Inches(0.15), step_w - Inches(0.4),
                 Inches(0.4), name, font=TITLE_FONT, size=13, bold=True,
                 color=color, line_spacing=1.0)
        add_text(s, x + Inches(0.2), step_y + Inches(0.55), step_w - Inches(0.4),
                 step_h - Inches(0.7), body, size=11, color=INK,
                 line_spacing=1.4)
        if i < 2:
            add_line(s, x + step_w, step_y + step_h / 2,
                     x + step_w + gap, step_y + step_h / 2,
                     color=GREY, weight=1.0)

    add_text(s, Inches(0.7), Inches(4.5), Inches(12), Inches(0.4),
             "IMPLEMENTATION NOTES",
             size=10, bold=True, color=TEAL, line_spacing=1.0)
    add_bullets(s, Inches(0.7), Inches(4.85), Inches(12), Inches(2.1), [
        "Per-country crosswalks live as YAML under etl/config/ — column renames, "
        "service-category mappings, and bounding-box validation.",
        "Facility UIDs are country-scoped and stable across reruns (TZA-000001, "
        "KEN-000001…). Downstream joins are safe.",
        "Outputs are plain CSV. The etl/ folder is a git submodule so both the "
        "working repo and the deploy clone consume the same version.",
    ], size=12)
    add_footer(s, page, total)


# ═══════════════════════════════════════════════════════════════════════════
# COMPONENT 1 — REGISTRY EXPLORER
# ═══════════════════════════════════════════════════════════════════════════

def slide_registry_what(total, page):
    s = new_slide()
    add_what_how_badge(s, Inches(0.7), Inches(0.4), "what")
    add_header(s, "Registry Explorer — coverage at a glance",
               eyebrow="Component 1 — what",
               subtitle="A single map + filter panel that loads 8 countries in parallel and stays snappy")

    add_bullets(s, Inches(0.7), Inches(2.55), Inches(7.5), Inches(3.5), [
        "Interactive world map with covered countries highlighted; hover for facility count.",
        "Client-side filters (country, facility type, ownership, name search) resolve "
        "instantly over the in-memory dataset.",
        "Filtered subset can be re-downloaded as CSV in one click.",
        "Per-country General Information panel appears below the map when a country is "
        "picked — type mix, ownership split, top regions.",
    ], size=14)

    x = Inches(8.6)
    add_text(s, x, Inches(2.55), Inches(4.2), Inches(0.4),
             "FACILITIES BY COUNTRY",
             size=11, bold=True, color=TEAL, line_spacing=1.0)
    counts = [
        ("Nigeria",  "51,023"),
        ("Ethiopia", "40,035"),
        ("Kenya",    "17,353"),
        ("Tanzania", "13,075"),
        ("Uganda",   "8,512"),
        ("Zambia",   "3,731"),
        ("Malawi",   "1,929"),
        ("Botswana", "1,076"),
    ]
    row_y = Inches(2.95)
    row_h = Inches(0.4)
    for i, (name, count) in enumerate(counts):
        y = row_y + row_h * i
        add_text(s, x + Inches(0.1), y, Inches(2.4), row_h, name,
                 size=12, color=INK, line_spacing=1.1)
        add_text(s, x + Inches(2.4), y, Inches(1.7), row_h, count,
                 size=12, bold=True, color=AMBER, align=PP_ALIGN.RIGHT,
                 line_spacing=1.1)
    add_footer(s, page, total)


def slide_registry_how(total, page):
    s = new_slide()
    add_what_how_badge(s, Inches(0.7), Inches(0.4), "how")
    add_header(s, "Client-side Leaflet + Papa Parse + Natural Earth GeoJSON",
               eyebrow="Component 1 — how",
               subtitle="No backend; the browser loads all 8 CSVs concurrently, then filters in JS")

    add_bullets(s, Inches(0.7), Inches(2.55), Inches(6.2), Inches(3.7), [
        "8 Papa Parse jobs launched in parallel, each streaming its country CSV.",
        "Country polygons drawn from Natural Earth via johan/world.geo.json.",
        "Filters (dropdowns + text search) apply against the in-memory Array once "
        "the last country loads.",
        "Download button rewrites the currently-filtered subset back to CSV in-browser "
        "and triggers a Blob download.",
        "All 8 country colour swatches unified to the HealthScope amber (#c9a227).",
    ], size=13)

    add_code_block(s, Inches(7.2), Inches(2.55), Inches(5.5), Inches(3.7),
                   "// Load all 8 countries concurrently\n"
                   "Object.entries(COUNTRIES).forEach(([k, cfg]) => {\n"
                   "  Papa.parse(cfg.file, {\n"
                   "    download: true, header: true,\n"
                   "    skipEmptyLines: true,\n"
                   "    complete: res => {\n"
                   "      res.data.forEach(r => {\n"
                   "        r._country = k;\n"
                   "        r._color   = cfg.color;\n"
                   "        allFacilities.push(r);\n"
                   "      });\n"
                   "      countryDone(k);\n"
                   "    }\n"
                   "  });\n"
                   "});",
                   size=11)
    add_footer(s, page, total)


# ═══════════════════════════════════════════════════════════════════════════
# COMPONENT 2 — POPULATION AND ACCESS
# ═══════════════════════════════════════════════════════════════════════════

def slide_pop_access_what(total, page):
    s = new_slide()
    add_what_how_badge(s, Inches(0.7), Inches(0.4), "what")
    add_header(s, "Population & Access — who each facility serves, and who can reach it",
               eyebrow="Component 2 — what",
               subtitle="Two side-by-side maps + right-side tables; toggles for facility points and protected areas")

    add_bullets(s, Inches(0.7), Inches(2.55), Inches(7.0), Inches(3.7), [
        "Left map: modelled travel-time to nearest facility, four coloured bands.",
        "Right map: Voronoi catchment polygons, coloured by population served.",
        "Right-side tables mirror the R outputs — national breakdown + worst-served "
        "districts + top-15 facilities by catchment.",
        "Facility points and WDPA protected areas are off by default; two check-boxes "
        "let the user layer them on.",
        "Tanzania is the flagship country — full analysis is running there; other "
        "seven countries in the queue.",
    ], size=13)

    # Right: TZ headline numbers
    x = Inches(8.3)
    add_text(s, x, Inches(2.55), Inches(4.3), Inches(0.4),
             "TANZANIA HEADLINE NUMBERS",
             size=11, bold=True, color=TEAL, line_spacing=1.0)
    tiles = [
        ("82.4 %", "within 30 min", GREEN_1),
        ("6.8 %",  ">90 min (4.74 M ppl)", RED),
        ("4,784",  "median Voronoi catchment", PURPLE),
        ("55 %",   "of Kaliua district >60 min", ORANGE),
    ]
    row_y = Inches(2.95)
    for i, (v, l, c) in enumerate(tiles):
        y = row_y + Inches(0.85) * i
        add_rect(s, x, y, Inches(4.3), Inches(0.72), fill=PALE, no_line=True)
        add_rect(s, x, y, Inches(0.08), Inches(0.72), fill=c, no_line=True)
        add_text(s, x + Inches(0.25), y + Inches(0.08), Inches(1.6),
                 Inches(0.55), v, font=TITLE_FONT, size=20, bold=True,
                 color=c, line_spacing=1.0)
        add_text(s, x + Inches(2.0), y + Inches(0.2), Inches(2.2),
                 Inches(0.5), l, size=11, color=INK, line_spacing=1.2)
    add_footer(s, page, total)


def slide_pop_access_how_method(total, page):
    s = new_slide()
    add_what_how_badge(s, Inches(0.7), Inches(0.4), "how")
    add_header(s, "Voronoi > radial buffers; travel time > straight-line distance",
               eyebrow="Component 2 — how (methodology)",
               subtitle="Two method choices that make the numbers honest")

    col_w = Inches(5.7)
    col_h = Inches(4.0)
    y = Inches(2.55)

    # Left: Voronoi
    add_rect(s, Inches(0.7), y, col_w, col_h, fill=WHITE,
             line_color=LIGHT_GREY, line_weight=1.0)
    add_rect(s, Inches(0.7), y, Inches(0.08), col_h, fill=PURPLE, no_line=True)
    add_text(s, Inches(0.95), y + Inches(0.15), col_w - Inches(0.3),
             Inches(0.4), "CATCHMENT METHOD — VORONOI",
             size=11, bold=True, color=PURPLE, line_spacing=1.0)
    add_text(s, Inches(0.95), y + Inches(0.55), col_w - Inches(0.3),
             Inches(0.5), "Nearest-facility polygons, no overlap",
             font=TITLE_FONT, size=15, bold=True, color=NAVY,
             line_spacing=1.1)
    add_bullets(s, Inches(0.95), y + Inches(1.15), col_w - Inches(0.3),
                Inches(2.6), [
        "Each point in the country belongs to exactly one facility.",
        "Total pop across all catchments ≈ national pop (no double-counting).",
        "Median catchment for Tanzania: 4,784 people/facility.",
        "Radial 5 km buffers overlap and inflate the median to ~12,700.",
    ], size=11)

    # Right: Travel time
    x2 = Inches(7.0)
    add_rect(s, x2, y, col_w, col_h, fill=WHITE,
             line_color=LIGHT_GREY, line_weight=1.0)
    add_rect(s, x2, y, Inches(0.08), col_h, fill=GREEN_1, no_line=True)
    add_text(s, x2 + Inches(0.25), y + Inches(0.15), col_w - Inches(0.3),
             Inches(0.4), "REACHABILITY — TRAVEL TIME BANDS",
             size=11, bold=True, color=GREEN_1, line_spacing=1.0)
    add_text(s, x2 + Inches(0.25), y + Inches(0.55), col_w - Inches(0.3),
             Inches(0.5), "Minutes-to-nearest, on a 100 m grid",
             font=TITLE_FONT, size=15, bold=True, color=NAVY,
             line_spacing=1.1)
    add_bullets(s, x2 + Inches(0.25), y + Inches(1.15), col_w - Inches(0.3),
                Inches(2.6), [
        "Modelled travel-time raster: minutes to reach nearest facility.",
        "Classified into 4 bands: 0-30, 30-60, 60-90, >90 min.",
        "Aggregated over WorldPop 2026 (100 m constrained) → population per band.",
        "Aggregated per district → ranked burden table.",
    ], size=11)

    add_rect(s, Inches(0.7), Inches(6.7), Inches(11.9), Inches(0.4),
             fill=PALE, no_line=True)
    add_text(s, Inches(0.9), Inches(6.75), Inches(11.5), Inches(0.3),
             "Buffers computed in Africa Albers Equal-Area (ESRI:102022) so distances stay accurate across multi-UTM-zone countries.",
             font=CODE_FONT, size=10, color=INK, line_spacing=1.1)
    add_footer(s, page, total)


def slide_pop_access_how_pipeline(total, page):
    s = new_slide()
    add_what_how_badge(s, Inches(0.7), Inches(0.4), "how")
    add_header(s, "4 R scripts, ~30 min end-to-end (Tanzania)",
               eyebrow="Component 2 — how (pipeline)",
               subtitle="Reproducible: raw inputs are gitignored; each script writes a portal-ready output")

    add_code_block(s, Inches(0.7), Inches(2.55), Inches(12), Inches(3.9),
                   "scripts/population/tanzania/\n"
                   "├─ 10_wdpa_mask.R          WDPA → strict IUCN categories (Ia, Ib, II, III)\n"
                   "│                           as optional map-overlay context.\n"
                   "│\n"
                   "├─ 11_travel_time_bands.R  Align 100 m modelled travel-time surface to\n"
                   "│                           100 m constrained WorldPop 2026, classify into\n"
                   "│                           4 bands, produce national + per-district rollups.\n"
                   "│\n"
                   "├─ 12_voronoi_catchments.R Voronoi polygons per facility, clipped to country,\n"
                   "│                           WorldPop summed inside each polygon.\n"
                   "│\n"
                   "└─ 13_simplify_bands.R     Post-process the 50 MB polygonised bands geojson\n"
                   "                            down to ~10 MB (1 km simplification tolerance).",
                   size=10)
    add_bullets(s, Inches(0.7), Inches(6.5), Inches(12), Inches(0.7), [
        "Windows R memory notes baked in: OpenBLAS thread cap, GDAL cache bump, "
        "terraOptions(todisk=TRUE), aggressive gc() between raster ops."
    ], size=11)
    add_footer(s, page, total)


# ═══════════════════════════════════════════════════════════════════════════
# COMPONENT 3 — SERVICE INTELLIGENCE
# ═══════════════════════════════════════════════════════════════════════════

def slide_services_what(total, page):
    s = new_slide()
    add_what_how_badge(s, Inches(0.7), Inches(0.4), "what")
    add_header(s, "Service Intelligence — what services exist where",
               eyebrow="Component 3 — what",
               subtitle="A per-country tour of the 7 canonical service domains and their per-facility profiles")

    add_bullets(s, Inches(0.7), Inches(2.55), Inches(7.0), Inches(3.7), [
        "Every facility record carries a per-domain service flag: malaria, HIV, "
        "maternity, laboratory, emergency, outpatient, inpatient.",
        "Country switcher shows service coverage: what percentage of facilities in "
        "the country offer each domain.",
        "Per-facility profile: click a facility, see the full list of services it "
        "runs (not just the top 5).",
        "Where the underlying MoH source doesn't publish services (e.g., Ethiopia), "
        "the coverage bar is greyed out rather than faked at zero.",
    ], size=13)

    # Right: 7 domains list
    x = Inches(8.3)
    add_text(s, x, Inches(2.55), Inches(4.3), Inches(0.4),
             "SEVEN CANONICAL DOMAINS",
             size=11, bold=True, color=TEAL, line_spacing=1.0)
    domains = [
        "Malaria services",
        "HIV / AIDS",
        "Maternity",
        "Laboratory",
        "Emergency",
        "Outpatient",
        "Inpatient",
    ]
    for i, d in enumerate(domains):
        y = Inches(2.95) + Inches(0.5) * i
        add_rect(s, x, y, Inches(0.4), Inches(0.4), fill=AMBER, no_line=True)
        add_text(s, x + Inches(0.1), y + Inches(0.06), Inches(0.3),
                 Inches(0.3), str(i + 1), font=TITLE_FONT, size=14, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER, line_spacing=1.0)
        add_text(s, x + Inches(0.55), y + Inches(0.08), Inches(3.7),
                 Inches(0.35), d, size=13, color=INK, line_spacing=1.0)
    add_footer(s, page, total)


def slide_services_how(total, page):
    s = new_slide()
    add_what_how_badge(s, Inches(0.7), Inches(0.4), "how")
    add_header(s, "Long-format service records aggregated in the browser",
               eyebrow="Component 3 — how",
               subtitle="~208k service rows for Tanzania alone — Papa Parse streams them, JS pivots them")

    add_bullets(s, Inches(0.7), Inches(2.55), Inches(6.2), Inches(3.7), [
        "Per-country service-records CSV: one row per (facility, service) pair.",
        "Papa Parse streams the file in chunks; UI renders progress as rows land.",
        "In-memory pivot maps every facility_code → { domain: [service_names…] }.",
        "Cross-country service comparison uses the crosswalk YAML from the ETL layer "
        "so 'Malaria RDT' in Tanzania collapses with 'Rapid malaria test' in Kenya.",
    ], size=13)

    add_code_block(s, Inches(7.2), Inches(2.55), Inches(5.5), Inches(3.7),
                   "// Group service rows by facility_code\n"
                   "const svcByCode = {};\n"
                   "svcRows.forEach(r => {\n"
                   "  const dom = domainOf(r.service_name);\n"
                   "  if (!dom) return;\n"
                   "  const s = svcByCode[r.facility_code]\n"
                   "        ?? (svcByCode[r.facility_code] = {});\n"
                   "  (s[dom] ??= []).push(r.service_name);\n"
                   "});\n"
                   "\n"
                   "// Now: svcByCode['100675-8'].malaria_services\n"
                   "//   = ['RDT test', 'ACT treatment', ...]",
                   size=11)
    add_footer(s, page, total)


# ═══════════════════════════════════════════════════════════════════════════
# COMPONENT 4 — FIND FACILITY
# ═══════════════════════════════════════════════════════════════════════════

def slide_find_facility_what(total, page):
    s = new_slide()
    add_what_how_badge(s, Inches(0.7), Inches(0.4), "what")
    add_header(s, "Find Facility — real-world routing from your location",
               eyebrow="Component 4 — what",
               subtitle="Public-facing search: pick a service, get the nearest N facilities by road, not by straight line")

    add_bullets(s, Inches(0.7), Inches(2.55), Inches(7.0), Inches(3.7), [
        "User's location detected via browser geolocation, or entered manually.",
        "Pick a service (malaria, maternity, laboratory…), pick a travel mode "
        "(drive / walk / cycle), pick how many results to show.",
        "Results ranked by real road distance and travel time — not straight-line.",
        "Facilities that fall off the road network gracefully drop back to "
        "straight-line and get a small 'straight-line' badge.",
    ], size=13)

    x = Inches(8.3)
    add_text(s, x, Inches(2.55), Inches(4.3), Inches(0.4),
             "TRAVEL MODES",
             size=11, bold=True, color=TEAL, line_spacing=1.0)
    modes = [
        ("Drive",  "driving-car",   "Road network, car speeds"),
        ("Walk",   "foot-walking",  "Roads + footpaths, ~5 km/h"),
        ("Cycle",  "cycling-regular", "Roads + cycleways, safer routes preferred"),
    ]
    for i, (label, key, sub) in enumerate(modes):
        y = Inches(2.95) + Inches(0.95) * i
        add_rect(s, x, y, Inches(4.3), Inches(0.82), fill=PALE, no_line=True)
        add_text(s, x + Inches(0.2), y + Inches(0.1), Inches(2.0),
                 Inches(0.35), label, font=TITLE_FONT, size=14, bold=True,
                 color=AMBER, line_spacing=1.0)
        add_text(s, x + Inches(0.2), y + Inches(0.45), Inches(4.0),
                 Inches(0.35), sub, size=10, color=INK, line_spacing=1.15)
        add_text(s, x + Inches(2.8), y + Inches(0.1), Inches(1.4),
                 Inches(0.35), key, font=CODE_FONT, size=10, color=GREY,
                 align=PP_ALIGN.RIGHT, line_spacing=1.0)
    add_footer(s, page, total)


def slide_find_facility_how(total, page):
    s = new_slide()
    add_what_how_badge(s, Inches(0.7), Inches(0.4), "how")
    add_header(s, "Pre-filter with Haversine, then one ORS Matrix call",
               eyebrow="Component 4 — how",
               subtitle="1 OpenRouteService request per search — quota-friendly, still gives real road distances")

    add_bullets(s, Inches(0.7), Inches(2.55), Inches(6.2), Inches(3.7), [
        "Client filters matching facilities by service, computes straight-line distance "
        "(Haversine) to every one.",
        "Take the 30 closest — these are the only candidates worth routing.",
        "Send those 30 destinations in a single ORS Matrix request; get back road "
        "distance + drive/walk/cycle time for each.",
        "Sort by real road distance, take top N (5, 10, 20 …), render on the map + list.",
        "API key lives in gitignored ors-key.js; referenced via a plain <script src> so "
        "it never appears in the qmd source.",
    ], size=13)

    add_code_block(s, Inches(7.2), Inches(2.55), Inches(5.5), Inches(3.7),
                   "// 1 request → up to 30 destinations\n"
                   "const resp = await fetch(\n"
                   "  'https://api.openrouteservice.org' +\n"
                   "  '/v2/matrix/' + currentProfile(),\n"
                   "  {\n"
                   "    method: 'POST',\n"
                   "    headers: {\n"
                   "      Authorization: ORS_API_KEY,\n"
                   "      'Content-Type': 'application/json'\n"
                   "    },\n"
                   "    body: JSON.stringify({\n"
                   "      locations,      // [origin, ...30 dests]\n"
                   "      sources: [0],\n"
                   "      destinations: [1..30],\n"
                   "      metrics: ['distance','duration'],\n"
                   "      units: 'km'\n"
                   "    })\n"
                   "  }\n"
                   ");",
                   size=10)
    add_footer(s, page, total)


# ═══════════════════════════════════════════════════════════════════════════
# COMPONENT 5 — MAPY CHATBOT
# ═══════════════════════════════════════════════════════════════════════════

def slide_mapy_what(total, page):
    s = new_slide()
    add_what_how_badge(s, Inches(0.7), Inches(0.4), "what")
    add_header(s, "Mapy — a scoped navigation assistant on every page",
               eyebrow="Component 5 — what",
               subtitle="Not a data oracle. A wayfinder for the portal + a spokesperson for the Malaria Atlas Project.")

    add_bullets(s, Inches(0.7), Inches(2.55), Inches(7.0), Inches(3.7), [
        "Floating chat bubble in the bottom-right corner of every page.",
        "Answers two kinds of questions: 'where do I find X in the portal?' and "
        "'what is the Malaria Atlas Project?'.",
        "Refuses (politely) to answer methodology questions, specific data queries, or "
        "medical advice — redirects the visitor to the relevant page instead.",
        "First message: \"Somebody has to see. Speak to Mapy.\" Gold theme (#c9a227), "
        "face avatar.",
    ], size=13)

    # Right column: guardrails
    x = Inches(8.3)
    add_text(s, x, Inches(2.55), Inches(4.3), Inches(0.4),
             "WHAT MAPY WON'T ANSWER",
             size=11, bold=True, color=TEAL, line_spacing=1.0)
    refuses = [
        ("×", "Specific data lookups\n(e.g. 'how many facilities in Kaliua?')"),
        ("×", "Deep methodology questions\n(Voronoi, WorldPop, ETL internals)"),
        ("×", "Medical advice or clinical decisions"),
        ("×", "Anything unrelated to HealthScope or MAP"),
    ]
    for i, (mark, txt) in enumerate(refuses):
        y = Inches(2.95) + Inches(0.85) * i
        add_rect(s, x, y, Inches(4.3), Inches(0.72), fill=PALE, no_line=True)
        add_text(s, x + Inches(0.15), y + Inches(0.15), Inches(0.4),
                 Inches(0.4), mark, font=TITLE_FONT, size=22, bold=True,
                 color=RED, align=PP_ALIGN.CENTER, line_spacing=1.0)
        add_text(s, x + Inches(0.65), y + Inches(0.1), Inches(3.5),
                 Inches(0.55), txt, size=10, color=INK, line_spacing=1.3)
    add_footer(s, page, total)


def slide_mapy_how(total, page):
    s = new_slide()
    add_what_how_badge(s, Inches(0.7), Inches(0.4), "how")
    add_header(s, "One Botpress embed, one system-prompt file",
               eyebrow="Component 5 — how",
               subtitle="Third-party webchat + a knowledge base seeded from the sitemap; portal-side integration is a single include")

    add_bullets(s, Inches(0.7), Inches(2.55), Inches(6.5), Inches(3.7), [
        "Botpress Cloud webchat v5 — free tier, floating-bubble launcher.",
        "Knowledge base crawled from map-data-engineering.github.io/healthscope-portal"
        "/sitemap.xml + MAP institutional pages.",
        "System prompt tightly scoped (see right) to prevent hallucinations on data or "
        "methodology.",
        "Portal-side embed: _chatbot-embed.html carries the two <script> tags, "
        "registered via _quarto.yml → include-after-body so every page inherits.",
        "Config lives on Botpress's CDN — dashboard changes propagate without a portal "
        "redeploy.",
    ], size=13)

    # Right column: system-prompt excerpt
    add_text(s, Inches(7.5), Inches(2.55), Inches(5.2), Inches(0.4),
             "SYSTEM PROMPT (EXCERPT)",
             size=11, bold=True, color=TEAL, line_spacing=1.0)
    add_code_block(s, Inches(7.5), Inches(2.95), Inches(5.2), Inches(3.5),
                   "You are HealthScope Assistant\n"
                   "(Mapy) — a navigation helper\n"
                   "for the HealthScope portal,\n"
                   "produced by the Malaria\n"
                   "Atlas Project.\n"
                   "\n"
                   "Your job is limited to two things:\n"
                   " 1. Portal navigation.\n"
                   " 2. MAP project context.\n"
                   "\n"
                   "Decline politely and redirect for:\n"
                   " - Specific data questions.\n"
                   " - Methodological details.\n"
                   " - Medical advice.\n"
                   "\n"
                   "Never invent facility counts\n"
                   "or district statistics.",
                   size=10)
    add_footer(s, page, total)


# ═══════════════════════════════════════════════════════════════════════════
# FINALIZE
# ═══════════════════════════════════════════════════════════════════════════

def slide_roadmap(total, page):
    s = new_slide()
    add_header(s, "What's shipped, what's next",
               eyebrow="Finalize — where we are, where we're heading",
               subtitle="Tanzania is the flagship; the same methodology will roll out to the other seven")

    y = Inches(2.55)
    col_h = Inches(4.4)
    col_w = Inches(5.9)

    # Shipped
    add_rect(s, Inches(0.7), y, col_w, col_h, fill=WHITE,
             line_color=GREEN_1, line_weight=1.25)
    add_rect(s, Inches(0.7), y, Inches(0.08), col_h, fill=GREEN_1,
             no_line=True)
    add_text(s, Inches(0.95), y + Inches(0.15), col_w - Inches(0.3),
             Inches(0.4), "SHIPPED",
             size=11, bold=True, color=GREEN_1, line_spacing=1.0)
    add_text(s, Inches(0.95), y + Inches(0.55), col_w - Inches(0.3),
             Inches(0.5), "Live on the public URL",
             font=TITLE_FONT, size=15, bold=True, color=NAVY,
             line_spacing=1.1)
    add_bullets(s, Inches(0.95), y + Inches(1.2), col_w - Inches(0.3),
                Inches(3.1), [
        "Registry Explorer across all 8 countries.",
        "Service Intelligence + Find Facility (client-side ORS routing).",
        "Population & Access — Tanzania flagship (Voronoi + travel-time + district ranking).",
        "Data Sources page with per-country MoH portal links.",
        "Mapy chatbot embedded on every page.",
    ], size=12)

    # Next
    x2 = Inches(6.8)
    add_rect(s, x2, y, col_w, col_h, fill=WHITE, line_color=AMBER,
             line_weight=1.25)
    add_rect(s, x2, y, Inches(0.08), col_h, fill=AMBER, no_line=True)
    add_text(s, x2 + Inches(0.25), y + Inches(0.15), col_w - Inches(0.3),
             Inches(0.4), "NEXT",
             size=11, bold=True, color=AMBER, line_spacing=1.0)
    add_text(s, x2 + Inches(0.25), y + Inches(0.55), col_w - Inches(0.3),
             Inches(0.5), "Blocking dependencies + follow-ups",
             font=TITLE_FONT, size=15, bold=True, color=NAVY,
             line_spacing=1.1)
    add_bullets(s, x2 + Inches(0.25), y + Inches(1.2), col_w - Inches(0.3),
                Inches(3.1), [
        "Travel-time surface per remaining country (7 countries).",
        "Constrained WorldPop 100 m per remaining country.",
        "Cloudflare Worker proxy for the ORS key (currently in gitignored client JS).",
        "Automated rerun of the R pipeline via GitHub Actions on data updates.",
        "Underserved-zones layer — needs the two items above for all countries first.",
    ], size=12)
    add_footer(s, page, total)


def slide_team(total, page):
    s = new_slide()
    add_header(s, "Team, sources, and how to get in touch",
               eyebrow="Finalize — acknowledgments")

    add_text(s, Inches(0.7), Inches(2.55), Inches(5.5), Inches(0.4),
             "TEAM", size=11, bold=True, color=TEAL, line_spacing=1.0)
    add_text(s, Inches(0.7), Inches(2.95), Inches(5.5), Inches(1.4),
             "Malaria Atlas Project — Dar es Salaam Node\n"
             "Data engineering, GIS, epidemiology, and portal build.",
             size=13, color=INK, line_spacing=1.55)

    add_text(s, Inches(0.7), Inches(4.35), Inches(5.5), Inches(0.4),
             "PRIMARY DATA SOURCES",
             size=11, bold=True, color=TEAL, line_spacing=1.0)
    add_bullets(s, Inches(0.7), Inches(4.75), Inches(5.5), Inches(2.3), [
        "MoH facility registries (per country — see the Data Sources page).",
        "WorldPop 2020 (1 km, unconstrained) + 2026 (100 m, constrained).",
        "OpenStreetMap + GeoFabrik extracts for road networks.",
        "WDPA (UNEP-WCMC & IUCN, May 2026) for protected-area masks.",
    ], size=12)

    add_text(s, Inches(7.0), Inches(2.55), Inches(5.6), Inches(0.4),
             "GET INVOLVED", size=11, bold=True, color=TEAL,
             line_spacing=1.0)
    add_text(s, Inches(7.0), Inches(3.0), Inches(5.6), Inches(3.5),
             "Live portal\n"
             "map-data-engineering.github.io/healthscope-portal\n"
             "\n"
             "Source (portal)\n"
             "github.com/map-data-engineering/healthscope-portal\n"
             "\n"
             "Source (ETL)\n"
             "github.com/map-data-engineering/health_facility_etl\n"
             "\n"
             "Malaria Atlas Project\n"
             "map.ox.ac.uk",
             size=12, color=INK, line_spacing=1.55)
    add_footer(s, page, total)


# ─── assemble the deck ─────────────────────────────────────────────────────
BUILDERS = [
    slide_title,                     # 1  Intro
    slide_what_is_healthscope,       # 2
    slide_the_gap,                   # 3
    slide_overview,                  # 4

    slide_foundation_what,           # 5  Foundation
    slide_foundation_how,            # 6

    slide_registry_what,             # 7  Component 1
    slide_registry_how,              # 8

    slide_pop_access_what,           # 9  Component 2
    slide_pop_access_how_method,     # 10
    slide_pop_access_how_pipeline,   # 11

    slide_services_what,             # 12 Component 3
    slide_services_how,              # 13

    slide_find_facility_what,        # 14 Component 4
    slide_find_facility_how,         # 15

    slide_mapy_what,                 # 16 Component 5
    slide_mapy_how,                  # 17

    slide_roadmap,                   # 18 Finalize
    slide_team,                      # 19
]
TOTAL = len(BUILDERS)

# Title slide takes no page number; the rest do.
BUILDERS[0](TOTAL)
for i, builder in enumerate(BUILDERS[1:], start=2):
    builder(TOTAL, i)

OUT = Path("HealthScope_Technical_Deck.pptx")
prs.save(OUT)
print(f"Wrote {OUT.resolve()}  ({TOTAL} slides)")
